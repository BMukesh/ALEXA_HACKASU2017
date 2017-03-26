import logging
import numpy as np

from random import randint

from flask import Flask, render_template

from flask_ask import Ask, statement, question, session

app = Flask(__name__)

ask = Ask(app, "/")

logging.getLogger("flask_ask").setLevel(logging.DEBUG)

@ask.launch

def new_game():
    welcome_msg = render_template('welcome')
    return question(welcome_msg)


@ask.intent("IntroIntent" , mapping={'action': 'Action'})

def next_round(action):
	print ("my action is ----> " + action)
	session.attributes['Action'] = action
	session.attributes['callFirst']=0
	session.attributes['wordP']=''
	session.attributes['type']='simple'
	session.attributes['difficult word']=[]
	session.attributes['attempt']=1
	session.attributes['slp']=1
	session.attributes['mlp']=1
	session.attributes['clp']=1
	session.attributes['rlp']=1
	session.attributes['dlp']=1
	return question('name of the file?')
	
@ask.intent("filenameIntent" , mapping={'filename': 'filename'})

def number_word(filename):
	session.attributes['filename'] = filename
	from pptx import Presentation
	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = "Hello, ASU!"
	subtitle.text = "We are creating a interactive projector with presentation skills!"
	prs.save(filename+'.pptx')
	return question('Okay. I am gona '+ session.attributes['Action']+', a presentation with file name '+ session.attributes['filename']+ '. What Next ?' )

@ask.intent("SpeechLearnIntent" , mapping={'word': 'Word'})

def speech_round(word):
	return question("please say the content you want to add in the next slide")

@ask.intent("CreateSlideIntent" , mapping={'title': 'title', "bullet":"bullet"})

def speech_round(title,bullet):
	from pptx import Presentation
	f = open(session.attributes['filename']+'.pptx', 'rb')
	prs = Presentation(f)
	bullet_slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	title_shape = shapes.title
	body_shape = shapes.placeholders[1]
	title_shape.text = title
	tf = body_shape.text_frame
	tf.text = bullet
	prs.save(session.attributes['filename']+'.pptx')
	f.close
	return statement("done")

if __name__ == '__main__':

    app.run(debug=True)
